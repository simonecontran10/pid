"""generator.py — orchestratore principale per export PPT dal PID."""
from __future__ import annotations

import hashlib
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from . import constants as C


# ============================================================
#  STEP 1 — SELEZIONE SLIDE
# ============================================================

def remove_slide_by_index(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId_list = list(sldIdLst)
    if index < 0 or index >= len(sldId_list):
        raise IndexError(f"Slide index {index} fuori range")
    sldId = sldId_list[index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    sldIdLst.remove(sldId)
    pres_rels = prs.part.rels
    if rId in pres_rels:
        pres_rels.pop(rId)


def keep_only_slides(prs, indices_to_keep):
    total = len(prs.slides)
    to_remove = sorted([i for i in range(total) if i not in indices_to_keep], reverse=True)
    for idx in to_remove:
        remove_slide_by_index(prs, idx)


# ============================================================
#  STEP 2 — SOSTITUZIONI TESTUALI
# ============================================================

def _iter_textboxes(shape):
    if shape.shape_type == 6:
        for child in shape.shapes:
            yield from _iter_textboxes(child)
    elif shape.has_text_frame:
        yield shape


def replace_text_in_slide(slide, replacements):
    count = 0
    for shape in slide.shapes:
        for tb in _iter_textboxes(shape):
            for para in tb.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if run.text == old:
                            run.text = new
                            count += 1
                            break
    return count


def _set_paragraph_text(paragraph, new_text):
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = new_text
    for run in runs[1:]:
        run._r.getparent().remove(run._r)


def fill_cover_header(cover_slide, stadio, data, competizione):
    for sh in cover_slide.shapes:
        for tb in _iter_textboxes(sh):
            text = tb.text_frame.text
            if "Stadio" in text and "Serie A" in text:
                paragraphs = tb.text_frame.paragraphs
                if len(paragraphs) >= 3:
                    _set_paragraph_text(paragraphs[0], stadio)
                    _set_paragraph_text(paragraphs[1], data)
                    _set_paragraph_text(paragraphs[2], competizione)
                return True
    return False


# ============================================================
#  STEP 3 — RIEMPIRE PLACEHOLDER CARD GIOCATORE
# ============================================================

def _short_surname(full_name):
    """Estrae il cognome da un full_name.
    - Default: ultima parola (es. "Federico Gatti" -> "Gatti")
    - Eccezione: particelle (Di, Van, De, etc.) -> particella + ultima (es. "Di Gregorio")
    - Eccezione: cognomi noti composti via lookup esplicito
    """
    if not full_name:
        return ""
    parts = full_name.strip().split()
    PARTICELLE = {"de", "del", "di", "da", "van", "von", "der", "den", "das", "dos", "la", "le"}
    # Cognomi noti composti — entrambe le parti sono cognome
    COMPOSTI_NOTI = {
        "joao mario": "Joao Mario",
        "thuram-ulien": "Thuram",
    }
    full_lower = full_name.lower()
    for k, v in COMPOSTI_NOTI.items():
        if k in full_lower:
            return v
    
    if len(parts) == 1:
        return parts[0]
    
    # Se la penultima è una particella, includi anche quella
    if len(parts) >= 2 and parts[-2].lower() in PARTICELLE:
        return f"{parts[-2].capitalize()} {parts[-1].capitalize()}"
    
    # Default: ultima parola
    return parts[-1]


def _foot_to_letter(foot):
    """Mappa il piede a sigla card: D=destro, S=sinistro, D/S=ambidestro"""
    if not foot:
        return "D"
    f = foot.lower().strip()
    if f == "left":
        return "S"
    if f == "both":
        return "D/S"
    return "D"


def _find_card_group(slide, card_code):
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        for child in sh.shapes:
            if child.has_text_frame:
                txt = child.text_frame.text.strip()
                if not txt:
                    continue
                first = txt.split()[0]
                if first == card_code:
                    return sh
                break
    return None


def fill_card_placeholder(slide, card_code, surname, height_cm, number, foot_letter, is_starter=True):
    grp = _find_card_group(slide, card_code)
    if grp is None:
        return False
    cognome_text = f"{surname.upper()} {height_cm} Cm" if height_cm else surname.upper()
    number_text = str(number) if number else ""
    foot_text = foot_letter

    for tb in _iter_textboxes(grp):
        if not tb.has_text_frame:
            continue
        para = tb.text_frame.paragraphs[0]
        if not para.runs:
            continue
        full_text = tb.text_frame.text.strip()
        first_word = full_text.split()[0] if full_text else ""
        if first_word == card_code or "Cm" in full_text:
            run = para.runs[0]
            run.text = cognome_text

    for child in grp.shapes:
        if child.shape_type != 6:
            continue
        for grandchild in child.shapes:
            if not grandchild.has_text_frame:
                continue
            txt = grandchild.text_frame.text.strip()
            if not txt:
                continue
            if txt.isdigit():
                grandchild.text_frame.paragraphs[0].runs[0].text = number_text
            elif txt in ("D", "S"):
                # foot_text può essere "D", "S", o "D/S" (ambidestro)
                tf = grandchild.text_frame
                para = tf.paragraphs[0]
                # Cancello tutti i run esistenti tranne il primo
                for r in para.runs[1:]:
                    r._r.getparent().remove(r._r)
                
                if foot_text == "S":
                    # Solo S, rosso bold
                    para.runs[0].text = "S"
                    para.runs[0].font.bold = True
                    para.runs[0].font.color.rgb = RGBColor.from_string("FF0000")
                elif foot_text == "D/S":
                    # D nero regolare + "/" + S rosso bold
                    para.runs[0].text = "D/"
                    para.runs[0].font.bold = False
                    para.runs[0].font.color.rgb = RGBColor.from_string("000000")
                    # Aggiungo run con "S" in rosso bold
                    from copy import deepcopy
                    new_r = deepcopy(para.runs[0]._r)
                    para.runs[0]._r.addnext(new_r)
                    # Imposto testo + style del nuovo run (devo ri-fetchare i runs)
                    new_run = para.runs[1]
                    new_run.text = "S"
                    new_run.font.bold = True
                    new_run.font.color.rgb = RGBColor.from_string("FF0000")
                else:
                    # D, nero regolare
                    para.runs[0].text = "D"
                    para.runs[0].font.bold = False
                    para.runs[0].font.color.rgb = RGBColor.from_string("000000")
    return True


def hide_card_placeholder(slide, card_code):
    grp = _find_card_group(slide, card_code)
    if grp is None:
        return False
    grp.left = Inches(-20)
    grp.top = Inches(-20)
    return True


def list_card_placeholders(slide):
    codes = []
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        for child in sh.shapes:
            if child.has_text_frame:
                txt = child.text_frame.text.strip()
                if not txt:
                    continue
                first = txt.split()[0]
                if re.match(r"^[A-Z]+\d+[A-Z]?$", first):
                    codes.append(first)
                    break
    return codes


# ============================================================
#  STEP 5 — RIEMPIRE SLIDE PRESSING
# ============================================================

PRESSING_LABEL_TO_PLACEHOLDER = {
    "Gk1":  "GK1",  "Rfb1": "RFB1", "Rcb1": "RCB1", "Cb1":  "CB1",
    "Lcb1": "LCB1", "Lfb1": "LFB1", "Rcm1": "RCM1", "Cm1":  "CM1",
    "Lcm1": "LCM1", "St1":  "ST1",  "St1b": "ST1B",
    "Rw1":  "RW1",  "Lw1":  "LW1",
}


def _find_pressing_group_by_label(slide, label):
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        for child in sh.shapes:
            if child.has_text_frame:
                if child.text_frame.text.strip() == label:
                    return sh, child
    return None, None


def _find_number_textbox_in_pressing_group(grp):
    for child in grp.shapes:
        if child.shape_type != 6:
            continue
        for grandchild in child.shapes:
            if grandchild.has_text_frame:
                txt = grandchild.text_frame.text.strip()
                if txt.isdigit():
                    return grandchild
    return None


def fill_pressing_slide(pressing_slide, players_dict):
    count = 0
    for press_label, formation_code in PRESSING_LABEL_TO_PLACEHOLDER.items():
        grp, label_tb = _find_pressing_group_by_label(pressing_slide, press_label)
        if grp is None:
            continue
        player_list = players_dict.get(formation_code)
        if not player_list:
            continue
        titolare = player_list[0]
        surname = _short_surname(titolare.get("name", ""))
        if surname:
            run = label_tb.text_frame.paragraphs[0].runs[0]
            run.text = " ".join(w.capitalize() for w in surname.split())
        number_tb = _find_number_textbox_in_pressing_group(grp)
        if number_tb:
            number = str(titolare.get("number", ""))
            if number:
                run = number_tb.text_frame.paragraphs[0].runs[0]
                run.text = number
        count += 1
    return count


# ============================================================
#  STEP 6 — SOSTITUZIONE LOGO SQUADRA
# ============================================================

LOGO_ITALIA_HASH = "af7418bf37"


def _replace_picture_blob(pic_shape, new_image_path, preserve_aspect=False, target_size_cm=None):
    """Sostituisce il blob immagine. 
    Se preserve_aspect=True, ridimensiona la shape per mantenere l'aspect ratio dell'immagine nuova.
    Se target_size_cm è specificato, usa quella dimensione (lato maggiore) invece dello spazio attuale.
    
    NOTA: gestisce correttamente il caso di image part CONDIVISE tra più PICTURE shape:
    crea una nuova image part dedicata invece di sovrascrivere il blob esistente
    (che cambierebbe anche le altre PICTURE che la condividono)."""
    NSMAP = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    blip_elements = pic_shape._element.findall(".//a:blip", NSMAP)
    if not blip_elements:
        return False
    blip = blip_elements[0]
    rId = blip.get(f"{{{NSMAP['r']}}}embed")
    if not rId:
        return False
    slide_part = pic_shape.part
    if rId not in slide_part.rels:
        return False
    
    with open(new_image_path, "rb") as f:
        new_blob = f.read()
    
    # Crea SEMPRE una nuova image part dedicata (evita modifiche a image part condivise)
    from pptx.parts.image import Image, ImagePart
    from pptx.opc.constants import CONTENT_TYPE as CT
    import io
    
    # Determine content type by extension
    suffix = Path(new_image_path).suffix.lower()
    if suffix in (".jpg", ".jpeg"):
        ct = CT.JPEG if hasattr(CT, "JPEG") else "image/jpeg"
    else:
        ct = CT.PNG if hasattr(CT, "PNG") else "image/png"
    
    # Crea una nuova Image e ImagePart
    new_image = Image.from_blob(new_blob, str(new_image_path))
    new_image_part = ImagePart.new(slide_part.package, new_image)
    new_rId = slide_part.relate_to(new_image_part, slide_part.rels[rId].reltype)
    
    # Aggiorna l'embed rId della shape
    blip.set(f"{{{NSMAP['r']}}}embed", new_rId)
    
    # Se richiesto, sistema le proporzioni (tipico per loghi quadrati in placeholder rettangolari)
    if preserve_aspect:
        try:
            from PIL import Image
            from io import BytesIO
            img = Image.open(BytesIO(new_blob))
            from pptx.util import Cm
            new_ratio = img.width / img.height
            cur_w = pic_shape.width
            cur_h = pic_shape.height
            cur_ratio = cur_w / cur_h
            
            # Centro nello spazio attuale (per il riposizionamento)
            center_x = pic_shape.left + cur_w // 2
            center_y = pic_shape.top + cur_h // 2
            
            if target_size_cm is not None:
                # Dimensione esplicita: lato maggiore = target_size_cm
                target_emu = Cm(target_size_cm)
                if new_ratio >= 1:
                    new_w = target_emu
                    new_h = int(target_emu / new_ratio)
                else:
                    new_h = target_emu
                    new_w = int(target_emu * new_ratio)
                pic_shape.width = new_w
                pic_shape.height = new_h
                pic_shape.left = center_x - new_w // 2
                pic_shape.top = center_y - new_h // 2
            elif abs(new_ratio - cur_ratio) > 0.05:
                # Adatta proporzioni al lato più piccolo dello spazio attuale
                if new_ratio > cur_ratio:
                    new_w = cur_w
                    new_h = int(cur_w / new_ratio)
                else:
                    new_h = cur_h
                    new_w = int(cur_h * new_ratio)
                pic_shape.width = new_w
                pic_shape.height = new_h
                pic_shape.left = center_x - new_w // 2
                pic_shape.top = center_y - new_h // 2
        except ImportError:
            pass  # Pillow non disponibile, lascia stretchato
        except Exception:
            pass  # In caso di errori, lascia stretchato
    return True


def replace_team_logos_all_slides(slides, new_logo_path):
    """Sostituisce il logo squadra in tutte le slide. Dimensioni:
    - Slide 0 (cover): 2.6 cm
    - Altre slide: 1.6 cm
    """
    count = 0
    targets = []  # lista di (pic_shape, slide_idx)

    def walk(sh):
        if sh.shape_type == 13:
            yield sh
        elif sh.shape_type == 6:
            for c in sh.shapes:
                yield from walk(c)

    for slide_idx, slide in enumerate(slides):
        for top_sh in slide.shapes:
            for pic in walk(top_sh):
                try:
                    h = hashlib.md5(pic.image.blob).hexdigest()[:10]
                except Exception:
                    h = None
                if h == LOGO_ITALIA_HASH:
                    targets.append((pic, slide_idx))

    for pic, slide_idx in targets:
        # Cover = slide 0 → 2.6 cm; altre slide → 1.6 cm
        # slide 0 = cover (2.6), slide 1 = formazione (1.6), slide 2 = pressing (1.1)
        target_cm = 2.6 if slide_idx == 0 else (1.6 if slide_idx == 1 else 1.1)
        if _replace_picture_blob(pic, new_logo_path, preserve_aspect=True, target_size_cm=target_cm):
            count += 1
    return count


def find_logo_path(team_name, loghi_dir):
    loghi_dir = Path(loghi_dir)
    candidates = [
        f"{team_name}.png",
        f"{team_name.lower()}.png",
        f"{team_name.upper()}.png",
        f"{team_name.title()}.png",
        f"{team_name.lower().replace(' ', '_')}.png",
        f"{team_name.lower().replace(' ', '-')}.png",
    ]
    for name in candidates:
        p = loghi_dir / name
        if p.exists():
            return p
    target = team_name.lower().replace(" ", "")
    for f in loghi_dir.iterdir():
        if f.suffix.lower() == ".png" and f.stem.lower().replace(" ", "").replace("_", "").replace("-", "") == target:
            return f
    return None


# ============================================================
#  STEP 7 — SOSTITUZIONE FOTO GIOCATORI
# ============================================================

# Mappatura PICTURE position → slot template, per ogni sistema
# La PICTURE è identificata dalla posizione (left, top) approssimativa.
# Tolerance ±0.15 inch.
PHOTO_POSITION_MAPS = {
    "3-5-2": {
        (4.70, 0.17): "GK1",
        (2.88, 1.38): "RCB1",
        (4.68, 1.39): "CB1",
        (6.53, 1.39): "LCB1",
        (7.67, 2.33): "LFB1",
        (1.65, 2.35): "RFB1",
        (4.71, 2.60): "CM1",
        (3.02, 3.06): "RCM1",
        (6.28, 3.06): "LCM1",
        (3.71, 4.29): "ST1",
        (5.69, 4.29): "ST1B",
    },
    "3-4-3": {
        (4.70, 0.17): "GK1",
        (2.88, 1.38): "RCB1",
        (4.68, 1.39): "CB1",
        (6.53, 1.39): "LCB1",
        (1.75, 2.61): "RFB1",
        (3.86, 2.61): "RCM1",
        (5.55, 2.61): "LCM1",
        (7.66, 2.61): "LFB1",
        (1.91, 3.97): "RW1",
        (7.39, 3.97): "LW1",
        (4.72, 4.29): "ST1",
    },
    "3-4-2-1": {
        (4.70, 0.17): "GK1",
        (2.88, 1.38): "RCB1",
        (4.68, 1.39): "CB1",
        (6.53, 1.39): "LCB1",
        (1.75, 2.61): "RFB1",
        (3.86, 2.61): "RCM1",
        (5.55, 2.61): "LCM1",
        (7.66, 2.61): "LFB1",
        (2.87, 3.76): "RW1",
        (6.58, 3.77): "LW1",
        (4.72, 4.29): "ST1",
    },
    "4-4-2": {
        (4.7, 0.17): "GK1",
        (1.83, 1.39): "RFB1",
        (3.69, 1.39): "RCB1",
        (5.65, 1.39): "LCB1",
        (7.6, 1.39): "LFB1",
        (3.76, 2.68): "RCM1",
        (5.65, 2.68): "LCM1",
        (1.91, 2.97): "RW1",
        (7.55, 2.97): "LW1",
        (3.71, 4.29): "ST1",
        (5.69, 4.29): "ST1B",
    },
    "4-3-3": {
        (4.7, 0.17): "GK1",
        (1.83, 1.39): "RFB1",
        (3.69, 1.39): "RCB1",
        (5.65, 1.39): "LCB1",
        (7.6, 1.39): "LFB1",
        (4.71, 2.6): "CM1",
        (3.06, 2.77): "RCM1",
        (6.32, 2.77): "LCM1",
        (1.91, 3.97): "RW1",
        (7.39, 3.97): "LW1",
        (4.72, 4.29): "ST1",
    },
    "4-2-3-1": {
        (4.7, 0.17): "GK1",
        (1.83, 1.08): "RFB1",
        (3.69, 1.08): "RCB1",
        (5.65, 1.08): "LCB1",
        (7.6, 1.08): "LFB1",
        (3.51, 2.22): "RCM1",
        (5.9, 2.22): "LCM1",
        (4.7, 3.24): "ST1B",
        (1.91, 3.37): "RW1",
        (7.55, 3.37): "LW1",
        (4.7, 4.4): "ST1",
    },
}

PHOTO_SIZE_INCHES = 0.59
POSITION_TOLERANCE = 0.15


def _find_player_photos(slide):
    """Trova le PICTURE 0.59x0.59 (foto giocatori) e ritorna lista (shape, left_in, top_in)."""
    photos = []
    for shape in slide.shapes:
        if shape.shape_type != 13:
            continue
        if not shape.width or not shape.height:
            continue
        w_in = shape.width.inches
        h_in = shape.height.inches
        # Filtra solo foto giocatori (0.59x0.59 ± tolleranza)
        if abs(w_in - PHOTO_SIZE_INCHES) > 0.05 or abs(h_in - PHOTO_SIZE_INCHES) > 0.05:
            continue
        l_in = shape.left.inches if shape.left else 0
        t_in = shape.top.inches if shape.top else 0
        photos.append((shape, round(l_in, 2), round(t_in, 2)))
    return photos


def _match_photo_to_slot(photo_pos, position_map):
    """Data una posizione (l, t), ritorna lo slot ID più vicino (entro tolerance)."""
    l, t = photo_pos
    best_slot = None
    best_dist = float("inf")
    for (mlmt, mt), slot_id in position_map.items():
        dist = abs(l - mlmt) + abs(t - mt)
        if dist < best_dist and dist <= POSITION_TOLERANCE * 2:
            best_dist = dist
            best_slot = slot_id
    return best_slot


def replace_player_photos(formation_slide, system, players, photos_dir):
    """Sostituisce le foto Sassuolo del template con quelle dei giocatori reali.
    
    Args:
        formation_slide: la slide della formazione (già selezionata)
        system: codice sistema, es. "3-5-2"
        players: dict del payload, es. {"GK1": [{"name": ..., "sots_id": 12345}], ...}
        photos_dir: Path alla cartella locale con le foto, es. data/photos/players_sots_lookup/
    
    Returns:
        dict con statistiche: {"replaced": N, "skipped_no_sots_id": N, "skipped_no_file": N, "skipped_no_match": N}
    """
    stats = {"replaced": 0, "skipped_no_sots_id": 0, "skipped_no_file": 0, "skipped_no_match": 0}
    photos_dir = Path(photos_dir)
    
    if system not in PHOTO_POSITION_MAPS:
        print(f"⚠️  PHOTO_POSITION_MAPS non definita per sistema '{system}'")
        return stats
    
    position_map = PHOTO_POSITION_MAPS[system]
    
    # Trova le 11 foto template
    template_photos = _find_player_photos(formation_slide)
    
    # Mappa shape → slot via posizioni
    slot_to_shape = {}
    for shape, l, t in template_photos:
        slot = _match_photo_to_slot((l, t), position_map)
        if slot:
            slot_to_shape[slot] = shape
    
    # Per ogni titolare nel payload, sostituisci la foto
    for slot_id, player_list in players.items():
        if not isinstance(player_list, list) or not player_list:
            continue
        titolare = player_list[0]
        sots_id = titolare.get("sots_id")
        if not sots_id:
            stats["skipped_no_sots_id"] += 1
            continue
        
        photo_file = photos_dir / f"{sots_id}.png"
        if not photo_file.exists():
            stats["skipped_no_file"] += 1
            continue
        
        target_shape = slot_to_shape.get(slot_id)
        if not target_shape:
            stats["skipped_no_match"] += 1
            continue
        
        # Sostituisce blob (riusa _replace_picture_blob già esistente)
        if _replace_picture_blob(target_shape, photo_file):
            stats["replaced"] += 1
    
    return stats




def bring_stripes_to_front(formation_slide):
    """Sposta strisce alto-sx in primo piano (Gruppo 9).
    Identificazione: GROUP shape con left e top entrambi negativi (off-canvas alto-sx)."""
    target = None
    for sh in formation_slide.shapes:
        if sh.shape_type != 6:
            continue
        if sh.left is not None and sh.top is not None:
            l_in = sh.left.inches
            t_in = sh.top.inches
            if l_in < -1 and t_in < -1:
                target = sh
                break
    if target is None:
        return False
    elem = target._element
    parent = elem.getparent()
    parent.remove(elem)
    parent.append(elem)  # in fondo all'XML = primo piano
    return True


def send_stripes_to_back(formation_slide):
    """Sposta strisce basso-dx dietro a TUTTO (Gruppo 13).
    Identificazione: GROUP shape con left positivo e top negativo, e size > 8 inches."""
    target = None
    for sh in formation_slide.shapes:
        if sh.shape_type != 6:
            continue
        if sh.left is not None and sh.top is not None and sh.width is not None:
            l_in = sh.left.inches
            t_in = sh.top.inches
            w_in = sh.width.inches
            # Strisce basso-dx: posizionate parzialmente off-canvas in alto-dx, dimensione grossa
            if l_in > 4 and t_in < 0 and w_in > 8:
                target = sh
                break
    if target is None:
        return False
    elem = target._element
    parent = elem.getparent()
    parent.remove(elem)
    # Inserisci all'inizio della spTree (= sotto a tutto)
    # Ma ATTENZIONE: i primi 2 elementi di spTree sono di solito nvGrpSpPr e grpSpPr (non shape vere)
    # Per essere sicuri, inserisco subito dopo grpSpPr
    NSMAP_P = "{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}"
    inserted = False
    children = list(parent)
    for i, child in enumerate(children):
        tag = child.tag
        # spTree contiene di solito: nvGrpSpPr, grpSpPr, poi le shape
        # Cerco la prima shape vera (sp, pic, grpSp, ecc.)
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}grpSp") or tag.endswith("}graphicFrame") or tag.endswith("}cxnSp"):
            child.addprevious(elem)
            inserted = True
            break
    if not inserted:
        # Fallback: insert all'inizio
        parent.insert(0, elem)
    return True





# ============================================================
#  STEP 8 — RICOLORAZIONE CERCHI PRESSING (Set A: squadra studiata)
# ============================================================

PRESSING_SET_A_LABELS = {"Gk1","Rfb1","Rcb1","Cb1","Lcb1","Lfb1","Rcm1","Cm1","Lcm1","St1","St1b","Rw1","Lw1"}


def _find_oval_in_group(grp):
    """Cerca primo Ovale (AUTO_SHAPE con name che contiene 'vale') nel gruppo."""
    for c in grp.shapes:
        if c.shape_type == 1 and hasattr(c, "name") and "vale" in (c.name or "").lower():
            return c
        if c.shape_type == 6:
            r = _find_oval_in_group(c)
            if r: return r
    return None


def _find_number_textbox_in_set_a_group(grp):
    """Cerca il TextBox del numero (testo numerico) dentro il gruppo Set A.
    Struttura: GROUP esterno → GROUP interno → [Ovale, TextBox numero] + TextBox label."""
    for c in grp.shapes:
        if c.shape_type == 6:
            for inner in c.shapes:
                if inner.has_text_frame:
                    t = inner.text_frame.text.strip()
                    if t and t.isdigit():
                        return inner
    return None


def recolor_pressing_set_a(pressing_slide, team_colors):
    """Ricolora i cerchi Set A (squadra studiata) della slide pressing.
    
    Args:
        pressing_slide: la slide pressing
        team_colors: dict {"fill": "RRGGBB", "border": "RRGGBB", "text": "RRGGBB"}
    
    Returns:
        numero di cerchi ricolorati (escluso Gk1 che resta nero per convenzione)
    """
    if not team_colors:
        return 0
    
    fill_hex = team_colors.get("fill", "FFFFFF")
    border_hex = team_colors.get("border", "000000")
    text_hex = team_colors.get("text", "000000")
    
    count = 0
    for sh in pressing_slide.shapes:
        if sh.shape_type != 6: continue
        # Trova label esterna del gruppo
        label = None
        for c in sh.shapes:
            if c.has_text_frame and c.shape_type != 6:
                t = c.text_frame.text.strip()
                if t and not t.isdigit():
                    label = t
                    break
        if label not in PRESSING_SET_A_LABELS:
            continue
        # Eccezione: Gk1 (portiere) resta nero per convenzione
        if label == "Gk1":
            continue
        
        oval = _find_oval_in_group(sh)
        if oval is None:
            continue
        
        # Cambio fill
        try:
            oval.fill.solid()
            oval.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
        except Exception:
            pass
        # Cambio border
        try:
            oval.line.color.rgb = RGBColor.from_string(border_hex)
        except Exception:
            pass
        
        # Cambio colore numero
        num_tb = _find_number_textbox_in_set_a_group(sh)
        if num_tb:
            for para in num_tb.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.color.rgb = RGBColor.from_string(text_hex)
                    except Exception:
                        pass
        
        count += 1
    return count


def load_team_colors(team_name, colori_dir):
    """Carica colori della squadra da assets/colori_squadre.json"""
    import json as _json
    colori_path = Path(colori_dir) / "colori_squadre.json"
    if not colori_path.exists():
        return None
    try:
        all_colors = _json.loads(colori_path.read_text())
    except Exception:
        return None
    key = (team_name or "").lower().strip()
    return all_colors.get(key) or all_colors.get("_default")



# ============================================================
#  ENTRY POINT
# ============================================================

def generate_pptx(payload, template_path, output_path, loghi_dir=None, photos_dir=None):
    team_name  = payload["team_name"]
    system     = payload["system"]
    match_info = payload.get("match_info") or {}
    players    = payload.get("players") or {}

    if system not in C.SYSTEM_TO_SLIDE_INDEX:
        raise ValueError(f"Sistema '{system}' non supportato")

    formation_slide_idx = C.SYSTEM_TO_SLIDE_INDEX[system]
    prs = Presentation(str(template_path))

    indices_to_keep = [C.COVER_SLIDE_INDEX, formation_slide_idx, C.PRESSING_SLIDE_INDEX]
    keep_only_slides(prs, indices_to_keep)

    cover_slide     = prs.slides[0]
    formation_slide = prs.slides[1]
    pressing_slide  = prs.slides[2]

    coach = match_info.get("coach", "")
    coach_parts = coach.strip().rsplit(" ", 1) if coach else ["", ""]
    coach_first = coach_parts[0] + " " if len(coach_parts) > 1 else ""
    coach_last  = coach_parts[-1] if coach_parts else ""

    cover_replacements = {"SQUADRA": team_name.upper()}
    cover_count = replace_text_in_slide(cover_slide, cover_replacements)
    if fill_cover_header(cover_slide,
                          stadio=match_info.get("stadio", ""),
                          data=match_info.get("data", ""),
                          competizione=match_info.get("competizione", "")):
        cover_count += 3

    formation_replacements = {
        "SQUADRA": team_name.upper(),
        "3-5-2": system,
        "Nome ": coach_first,
        "Cognome": coach_last,
    }
    formation_count = replace_text_in_slide(formation_slide, formation_replacements)

    # Allarga il textbox COACH a sx per evitare a capo
    for sh in formation_slide.shapes:
        if sh.has_text_frame and "OACH" in sh.text_frame.text:
            from pptx.util import Inches
            current_left = sh.left
            current_width = sh.width
            new_width_emu = Inches(3.5)
            if new_width_emu > current_width:
                new_left = current_left + current_width - new_width_emu
                sh.left = new_left
                sh.width = new_width_emu
            break

    # Allarga il textbox COACH a sx per evitare a capo
    for sh in formation_slide.shapes:
        if sh.has_text_frame and "OACH" in sh.text_frame.text:
            from pptx.util import Inches
            current_left = sh.left
            current_width = sh.width
            new_width_emu = Inches(3.5)
            if new_width_emu > current_width:
                # Espande a sinistra (mantenendo il bordo dx)
                new_left = current_left + current_width - new_width_emu
                sh.left = new_left
                sh.width = new_width_emu
            break

    all_codes = list_card_placeholders(formation_slide)
    used_codes = set()
    fill_stats = {"filled": 0, "skipped_extra": 0, "skipped_missing": 0}

    for slot_key, player_list in players.items():
        if not isinstance(player_list, list):
            continue
        m = re.match(r"^([A-Z]+)(\d+)([A-Z]?)$", slot_key)
        if not m:
            continue
        prefix = m.group(1)
        num_base = int(m.group(2))
        letter_suffix = m.group(3)
        for i, player in enumerate(player_list):
            if i > 4:
                fill_stats["skipped_extra"] += 1
                continue
            target_num = num_base + i
            target_code = f"{prefix}{target_num}{letter_suffix}"
            if target_code not in all_codes:
                fill_stats["skipped_missing"] += 1
                continue
            surname = _short_surname(player.get("name", ""))
            height = player.get("height", 0) or 0
            number = player.get("number", "")
            foot = _foot_to_letter(player.get("foot"))
            ok = fill_card_placeholder(formation_slide, target_code,
                                       surname=surname, height_cm=height,
                                       number=number, foot_letter=foot,
                                       is_starter=(i == 0))
            if ok:
                fill_stats["filled"] += 1
                used_codes.add(target_code)

    hidden_count = 0
    for code in all_codes:
        if code not in used_codes:
            if hide_card_placeholder(formation_slide, code):
                hidden_count += 1

    # Ricolora cerchi Set A PRIMA di fill_pressing_slide (sennò le label vengono sostituite con cognomi)
    pressing_recolored = 0
    if loghi_dir:
        team_colors = load_team_colors(team_name, loghi_dir.parent if hasattr(loghi_dir, 'parent') else loghi_dir)
        if team_colors:
            pressing_recolored = recolor_pressing_set_a(pressing_slide, team_colors)

    pressing_filled = fill_pressing_slide(pressing_slide, players)

    # Sostituisce foto giocatori del template con quelle reali
    photos_stats = {"replaced": 0}
    photos_dir_path = None
    if photos_dir:
        photos_stats = replace_player_photos(formation_slide, system, players, photos_dir)
        photos_dir_path = str(photos_dir)

    # Sposta strisce alto-sx in primo piano (sopra logo squadra)
    bring_stripes_to_front(formation_slide)
    # Sposta strisce basso-dx in fondo (dietro sfondo campo)
    send_stripes_to_back(formation_slide)

    logo_replaced = 0
    logo_path = None
    if loghi_dir:
        logo_path = find_logo_path(team_name, loghi_dir)
        if logo_path:
            logo_replaced = replace_team_logos_all_slides(
                [cover_slide, formation_slide, pressing_slide], logo_path)
        else:
            print(f"⚠️  Logo per '{team_name}' non trovato in {loghi_dir}")

    prs.save(str(output_path))

    return {
        "cover_replacements": cover_count,
        "formation_replacements": formation_count,
        "card_filled": fill_stats["filled"],
        "card_hidden": hidden_count,
        "card_skipped_extra": fill_stats["skipped_extra"],
        "card_skipped_missing": fill_stats["skipped_missing"],
        "pressing_filled": pressing_filled,
        "pressing_recolored": pressing_recolored,
        "logo_replaced": logo_replaced,
        "logo_path": str(logo_path) if logo_path else None,
        "output": str(output_path),
        "photos_replaced": photos_stats.get("replaced", 0),
        "photos_skipped_no_sots_id": photos_stats.get("skipped_no_sots_id", 0),
        "photos_skipped_no_file": photos_stats.get("skipped_no_file", 0),
        "photos_skipped_no_match": photos_stats.get("skipped_no_match", 0),
    }
